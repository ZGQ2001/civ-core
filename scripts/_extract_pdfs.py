"""批量提取培训目录下所有 PDF/PPTX 文本到 data/training_materials/"""
import os, sys
from pathlib import Path

SRC = Path(r"G:\我的云端硬盘\工作\培训")
DST = Path(r"D:\CodeProjects\civ-core\data\training_materials")
DST.mkdir(parents=True, exist_ok=True)

# ---- PDF 提取 ----
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# ---- PPTX 提取 ----
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

for fpath in sorted(SRC.iterdir()):
    if fpath.is_dir():
        continue
    name = fpath.name
    print(f"Processing: {name}")
    
    text = ""
    if fpath.suffix.lower() == ".pdf" and fitz:
        doc = fitz.open(str(fpath))
        for page in doc:
            text += page.get_text() + "\n---PAGE_BREAK---\n"
        doc.close()
    elif fpath.suffix.lower() == ".pptx" and Presentation:
        prs = Presentation(str(fpath))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text += para.text + "\n"
            text += "---SLIDE_BREAK---\n"
    else:
        print(f"  SKIP: no extractor for {fpath.suffix}")
        continue
    
    out_path = DST / f"{fpath.stem}.txt"
    out_path.write_text(text, encoding="utf-8")
    print(f"  -> {out_path} ({len(text)} chars)")

print("\nDone.")
